from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT.parent / "ppt_template_proj2.pptx"
OUTPUT = ROOT / "proj2_kst.pptx"

BLUE = RGBColor(0x1C, 0x73, 0xFF)
CYAN = RGBColor(0x50, 0xC4, 0xEE)
NAVY = RGBColor(0x12, 0x2C, 0x59)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x59, 0x59, 0x59)
MID_GRAY = RGBColor(0x98, 0xA2, 0xB3)
LIGHT_GRAY = RGBColor(0xF4, 0xF7, 0xFC)
BAR_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
PALE_BLUE = RGBColor(0xE9, 0xF2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE = "Gwangyang Sunshine Bold"
FONT_BODY = "Pretendard Regular"
FONT_BOLD = "Pretendard SemiBold"


def remove_all_slides(prs):
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def no_line(shape):
    shape.line.fill.background()


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_shape(slide, kind, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    if line is None:
        no_line(shape)
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.4)
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=22,
    color=BLACK,
    font=FONT_BODY,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.04,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_rich_title(slide, y, segments, size=54):
    box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(23.67), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for text, color in segments:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_TITLE
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_picture_contain(slide, path, x, y, w, h):
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw = iw * scale
    ph = ih * scale
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_base(slide, section, deck_name="알약 각인 OCR", draw_header=True):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 26.667, 15, BLUE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, -0.45, 25.23, 14.72, WHITE)
    if draw_header:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.68, 0.98, 23.72, 0.86, BAR_GRAY)
        add_text(slide, 2.3, 1.08, 7.5, 0.55, section, 19, BLACK)
        add_text(slide, 18.4, 1.08, 6.35, 0.55, deck_name, 18, BLACK, align=PP_ALIGN.RIGHT)


def add_header_foreground(slide, section, deck_name="알약 각인 OCR"):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.68, 0.98, 23.72, 0.86, BAR_GRAY)
    number, label = section.split(" ", 1)
    add_text(slide, 2.3, 1.08, 0.82, 0.55, number, 19, BLACK)
    add_text(slide, 3.18, 1.08, 4.5, 0.55, label, 19, BLACK)
    add_text(slide, 21.15, 1.08, 3.60, 0.55, deck_name, 17, BLACK, align=PP_ALIGN.CENTER)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_arrow(slide, x, y, w=0.52, h=0.34, color=BLUE):
    return add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x, y, w, h, color)


def add_issue_card(slide, x, title, chip, description, image_paths):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.35, 7.42, 7.95, LIGHT_GRAY)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.45, 4.67, 6.52, 0.78, BLUE)
    add_text(
        slide,
        x + 0.58,
        4.75,
        6.26,
        0.58,
        title,
        24,
        WHITE,
        FONT_BOLD,
        True,
        PP_ALIGN.CENTER,
    )

    if len(image_paths) == 2:
        image_w = 3.02
        for idx, path in enumerate(image_paths):
            image_x = x + 0.49 + idx * 3.42
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, image_x, 5.72, image_w, 2.42, WHITE)
            add_picture_contain(slide, path, image_x + 0.08, 5.81, image_w - 0.16, 2.24)
    else:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.49, 5.72, 6.44, 3.28, WHITE)
        add_picture_contain(slide, image_paths[0], x + 0.68, 5.86, 6.06, 3.00)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 1.24, 9.27, 4.94, 0.70, WHITE, BLUE)
    add_text(
        slide,
        x + 1.37,
        9.34,
        4.68,
        0.52,
        chip,
        18,
        BLUE,
        FONT_BOLD,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x + 0.62,
        10.18,
        6.18,
        1.45,
        description,
        18,
        DARK_GRAY,
        FONT_BODY,
        False,
        PP_ALIGN.CENTER,
        line_spacing=1.08,
    )


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, "01. 문제 정의")
    add_rich_title(slide, 2.16, [("알약 각인은 항상 ", BLACK), ("한 줄", BLUE), ("로 읽히지 않습니다", BLACK)])
    add_text(
        slide,
        2.7,
        3.18,
        21.27,
        0.62,
        "2줄 배치, 꽉 찬 각인, 유사 문자 혼동은 순차 인식에 불리한 조건입니다.",
        23,
        DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_issue_card(
        slide,
        1.65,
        "2줄 각인",
        "행 구분 필요",
        "문자열을 한 줄로 펼치면\n읽기 순서가 모호해집니다.",
        [ROOT / "ocr_0_bg.png", ROOT / "ocr_1_bg.png"],
    )
    add_issue_card(
        slide,
        9.62,
        "알약을 가득 채운 각인",
        "문자 경계 약화",
        "여백이 부족하면 문자와 외곽을\n안정적으로 분리하기 어렵습니다.",
        [ROOT / "ocr_2_bg.png"],
    )
    add_issue_card(
        slide,
        17.59,
        "숫자·알파벳 혼동",
        "S ↔ 5  /  O ↔ 0",
        "형태가 비슷한 문자는 주변 문맥이\n부족할 때 혼동되기 쉽습니다.",
        [ROOT / "ocr_3_bg.png"],
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.65, 12.58, 23.36, 0.92, NAVY)
    add_text(
        slide,
        2.10,
        12.67,
        22.46,
        0.68,
        "한 줄 시퀀스 가정만으로는 알약 각인의 2차원 배치와 전역 문맥을 충분히 표현하기 어렵습니다.",
        20,
        WHITE,
        FONT_BOLD,
        True,
        PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "기존 CRNN은 규칙적인 한 줄 텍스트에는 효율적이지만 실제 알약 각인은 항상 한 줄로 배치되지 않습니다. 두 줄 각인은 읽기 순서가 모호하고, 알약 면을 가득 채운 각인은 문자 경계를 안정적으로 분리하기 어렵습니다. 또한 S와 5, O와 0처럼 형태가 비슷한 문자는 주변 문맥이 부족할 때 혼동되기 쉽습니다. 따라서 2차원 배치와 전체 영역의 관계를 함께 보는 구조가 필요합니다.",
    )


def add_table_cell(slide, x, y, w, h, text, fill, color=BLACK, bold=False, align=PP_ALIGN.LEFT, size=18, line=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line)
    return add_text(
        slide,
        x + 0.18,
        y + 0.08,
        w - 0.36,
        h - 0.16,
        text,
        size,
        color,
        FONT_BOLD if bold else FONT_BODY,
        bold,
        align,
        margin=0.02,
        line_spacing=1.02,
    )


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, "02. 모델 비교", draw_header=False)
    add_rich_title(slide, 2.16, [("CRNN에서 ", BLACK), ("공간 이해형 OCR", BLUE), ("로", BLACK)])
    add_text(
        slide,
        2.7,
        3.18,
        21.27,
        0.62,
        "효율적인 순차 인식과 2차원 전역 문맥 인식의 차이를 비교했습니다.",
        23,
        DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

    x0, x1, x2 = 1.65, 5.65, 14.85
    w0, w1, w2 = 4.00, 9.20, 10.16
    y = 4.34
    header_h = 0.95
    add_table_cell(slide, x0, y, w0, header_h, "비교 항목", NAVY, WHITE, True, PP_ALIGN.CENTER, 21)
    add_table_cell(slide, x1, y, w1, header_h, "CRNN", MID_GRAY, WHITE, True, PP_ALIGN.CENTER, 22)
    add_table_cell(slide, x2, y, w2, header_h, "CNN + Transformer", BLUE, WHITE, True, PP_ALIGN.CENTER, 22)

    rows = [
        ("주요 강점", "규칙적인 한 줄 텍스트를\n효율적으로 인식", "각인 영역 전체의\n전역 문맥 파악"),
        ("공간 정보 처리", "특징 맵을 1차원\n시퀀스로 변환", "2D 위치 인코딩으로\n행과 열 정보 유지"),
        ("어려운 배치", "문자 간격이 넓거나 정렬이\n불규칙하면 대응이 제한적", "어텐션으로 멀리 떨어지거나\n불규칙한 문자를 연결"),
        ("연산량", "비교적 적음", "시각 토큰 수가 증가할수록\n연산량 증가"),
    ]
    row_h = 1.62
    for idx, (criterion, crnn, transformer) in enumerate(rows):
        row_y = y + header_h + idx * row_h
        criterion_fill = PALE_BLUE if idx % 2 == 0 else LIGHT_GRAY
        body_fill = WHITE if idx % 2 == 0 else LIGHT_GRAY
        add_table_cell(slide, x0, row_y, w0, row_h, criterion, criterion_fill, NAVY, True, PP_ALIGN.CENTER, 19)
        add_table_cell(slide, x1, row_y, w1, row_h, crnn, body_fill, DARK_GRAY, False, PP_ALIGN.CENTER, 18)
        add_table_cell(slide, x2, row_y, w2, row_h, transformer, body_fill, BLACK, True, PP_ALIGN.CENTER, 18)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.65, 12.10, 23.36, 1.28, BLUE)
    add_text(
        slide,
        2.20,
        12.24,
        22.26,
        0.98,
        "핵심 메시지  |  불규칙한 각인 배치를 더 잘 찾기 위해 더 높은 연산 비용을 감수합니다.",
        21,
        WHITE,
        FONT_BOLD,
        True,
        PP_ALIGN.CENTER,
    )
    add_header_foreground(slide, "02. 모델 비교")
    add_notes(
        slide,
        "CRNN은 규칙적인 한 줄 텍스트를 적은 연산량으로 인식하는 데 강점이 있습니다. 하지만 특징 맵을 1차원 시퀀스로 바꾸는 과정에서 행과 열의 위치 정보가 약해질 수 있습니다. CNN과 Transformer를 결합한 구조는 2차원 위치 정보를 유지하고 어텐션으로 멀리 떨어진 문자 관계를 연결할 수 있습니다. 대신 시각 토큰 수가 늘어날수록 연산 비용이 증가합니다. 이 프로젝트는 실제 알약의 불규칙한 각인 배치를 더 잘 표현하기 위해 이 비용을 감수합니다.",
    )


def add_flow_node(slide, x, y, w, h, number, label, fill, text_color=WHITE):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.12, y + 0.34, 0.55, 0.55, CYAN if fill != CYAN else NAVY)
    add_text(slide, x + 0.12, y + 0.34, 0.55, 0.55, number, 13, WHITE, FONT_BOLD, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.75, y + 0.15, w - 0.88, h - 0.30, label, 17, text_color, FONT_BOLD, True, PP_ALIGN.CENTER)


def add_sequence_visual(slide, x, y):
    labels = ["P", "D", "1", "0", "0"]
    for idx, label in enumerate(labels):
        box_x = x + idx * 0.78
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, box_x, y, 0.62, 0.62, WHITE, BLUE)
        add_text(slide, box_x, y, 0.62, 0.62, label, 17, BLUE, FONT_BOLD, True, PP_ALIGN.CENTER)
        if idx < len(labels) - 1:
            add_arrow(slide, box_x + 0.64, y + 0.20, 0.13, 0.20, MID_GRAY)


def add_attention_visual(slide, x, y):
    points = []
    for row in range(3):
        for col in range(5):
            points.append((x + col * 0.78, y + row * 0.42))
    links = [(0, 8), (1, 13), (2, 9), (4, 10), (5, 14), (7, 12), (3, 11)]
    for start, end in links:
        x1, y1 = points[start]
        x2, y2 = points[end]
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1 + 0.15),
            Inches(y1 + 0.15),
            Inches(x2 + 0.15),
            Inches(y2 + 0.15),
        )
        line.line.color.rgb = CYAN
        line.line.width = Pt(1.2)
        line.line.transparency = 30
    for idx, (px, py) in enumerate(points):
        fill = BLUE if idx in (0, 4, 7, 11, 14) else WHITE
        line_color = BLUE if fill == WHITE else None
        add_shape(slide, MSO_SHAPE.OVAL, px, py, 0.30, 0.30, fill, line_color)


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, "03. 모델 구조")
    add_rich_title(slide, 2.16, [("순차적 읽기에서 ", BLACK), ("공간적 이해", BLUE), ("로", BLACK)])
    add_text(
        slide,
        2.7,
        3.18,
        21.27,
        0.62,
        "CRNN은 읽기 순서에 집중하고, 프로젝트 모델은 2차원 위치와 전체 문맥을 함께 봅니다.",
        23,
        DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

    panels = [(1.65, "CRNN", MID_GRAY), (13.67, "CNN + Transformer", BLUE)]
    for x, title, color in panels:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.30, 11.34, 8.98, LIGHT_GRAY)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.46, 4.62, 10.42, 0.82, color)
        add_text(slide, x + 0.62, 4.70, 10.10, 0.62, title, 25, WHITE, FONT_BOLD, True, PP_ALIGN.CENTER)

    left_nodes = [
        (2.14, "01", "CNN\n특징"),
        (5.70, "02", "1D\n시퀀스"),
        (9.26, "03", "순환 신경망\n인식"),
    ]
    for idx, (x, number, label) in enumerate(left_nodes):
        add_flow_node(slide, x, 5.88, 2.70, 1.28, number, label, MID_GRAY)
        if idx < len(left_nodes) - 1:
            add_arrow(slide, x + 2.82, 6.34, 0.55, 0.34, MID_GRAY)

    right_nodes = [
        (14.11, "01", "CNN\n특징"),
        (16.79, "02", "2D 위치\n정보"),
        (19.47, "03", "전역\n어텐션"),
        (22.15, "04", "각인\n디코딩"),
    ]
    for idx, (x, number, label) in enumerate(right_nodes):
        add_flow_node(slide, x, 5.88, 2.10, 1.28, number, label, BLUE)
        if idx < len(right_nodes) - 1:
            add_arrow(slide, x + 2.20, 6.34, 0.38, 0.34, BLUE)

    add_text(slide, 2.32, 7.58, 10.00, 0.50, "가로 방향으로 펼친 뒤 순서대로 읽기", 18, DARK_GRAY, FONT_BOLD, True, PP_ALIGN.CENTER)
    add_sequence_visual(slide, 5.30, 8.30)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.28, 9.45, 10.08, 2.32, WHITE)
    add_text(
        slide,
        2.68,
        9.74,
        9.28,
        1.72,
        "CNN이 추출한 특징을 가로 방향의 시퀀스로 변환하여,\n글자를 정해진 순서에 따라 차례대로 인식합니다.",
        19,
        DARK_GRAY,
        FONT_BODY,
        False,
        PP_ALIGN.CENTER,
        line_spacing=1.08,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.18, 12.06, 6.28, 0.66, WHITE, MID_GRAY)
    add_text(slide, 4.35, 12.12, 5.94, 0.50, "강점  |  빠르고 효율적인 순차 인식", 17, DARK_GRAY, FONT_BOLD, True, PP_ALIGN.CENTER)

    add_text(slide, 14.35, 7.58, 9.98, 0.50, "행·열 위치를 유지하며 전체 관계 연결", 18, BLUE, FONT_BOLD, True, PP_ALIGN.CENTER)
    add_attention_visual(slide, 17.55, 8.18)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 14.30, 9.45, 10.08, 2.32, WHITE)
    add_text(
        slide,
        14.70,
        9.74,
        9.28,
        1.72,
        "2차원 위치 정보를 유지한 채 전체 영역의 관계를 동시에 분석하여,\n다양한 방향과 배열의 각인을 효과적으로 인식합니다.",
        19,
        DARK_GRAY,
        FONT_BODY,
        False,
        PP_ALIGN.CENTER,
        line_spacing=1.08,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 16.20, 12.06, 6.28, 0.66, WHITE, BLUE)
    add_text(slide, 16.37, 12.12, 5.94, 0.50, "강점  |  불규칙한 2D 배치 이해", 17, BLUE, FONT_BOLD, True, PP_ALIGN.CENTER)

    add_notes(
        slide,
        "CRNN은 CNN 특징을 가로 방향의 시퀀스로 바꾸고 순환 신경망이 왼쪽에서 오른쪽으로 차례대로 읽습니다. 규칙적인 한 줄 텍스트에는 효율적이지만 2차원 배치를 충분히 보존하기 어렵습니다. 프로젝트 모델은 CNN 특징에 2차원 위치 정보를 더하고 전역 어텐션으로 전체 영역의 관계를 동시에 분석합니다. 이를 통해 서로 다른 행이나 멀리 떨어진 위치의 문자도 연결해 각인을 디코딩할 수 있습니다.",
    )


def main():
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)
    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    prs.core_properties.title = "CNN + Transformer를 이용한 알약 각인 인식"
    prs.core_properties.subject = "CRNN의 한계와 CNN + Transformer 기반 공간적 각인 인식"
    prs.core_properties.author = "프로젝트 2팀"
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
